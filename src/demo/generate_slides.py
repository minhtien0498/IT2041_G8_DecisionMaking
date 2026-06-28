"""
Script to generate PowerPoint slides for the Midterm presentation.
Requires: python-pptx (pip install python-pptx)
Output: outputs/midterm_presentation.pptx
"""

import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: python-pptx library not found.")
    print("Please install it by running: pip install python-pptx")
    import sys
    sys.exit(1)

# Initialize presentation
prs = Presentation()

# Set slide size to widescreen (16:9)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette (Premium Light Theme)
BG_COLOR = RGBColor(248, 250, 252)      # Slate 50 (Off-white background)
TEXT_MAIN = RGBColor(15, 23, 42)        # Slate 900 (Dark Slate main text)
TEXT_MUTED = RGBColor(100, 116, 139)    # Slate 500 (Cool Gray muted text)
ACCENT_BLUE = RGBColor(29, 78, 216)     # Blue 700 (Rich Blue accent)
ACCENT_GOLD = RGBColor(180, 83, 9)      # Amber 700 (Deep Amber accent)
CARD_BG = RGBColor(255, 255, 255)       # Pure White card background

def set_slide_background(slide):
    """Set slide background to solid light Slate 50."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_title(slide, text, subtitle_text=None):
    """Add a consistent slide header."""
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Arial'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.name = 'Arial'
        p2.font.size = Pt(16)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(8)

def create_card(slide, left, top, width, height, title, items, title_color=ACCENT_BLUE):
    """Create a structured information card with custom styled bullet points."""
    # Add background rectangle for the card
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = RGBColor(203, 213, 225) # Slate 300 line border for light card
    shape.line.width = Pt(1.5)
    
    # Add text inside the card
    tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.25), width - Inches(0.5), height - Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    # Card Title
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = 'Arial'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.space_after = Pt(12)
    
    # Card Bullets
    for item in items:
        p_item = tf.add_paragraph()
        p_item.text = item
        p_item.font.name = 'Arial'
        p_item.font.size = Pt(14)
        p_item.font.color.rgb = TEXT_MAIN
        p_item.space_after = Pt(8)
        p_item.level = 0

# ==============================================================================
# SLIDE 1: Title Slide (Cover)
# ==============================================================================
slide_layout = prs.slide_layouts[6] # Blank layout
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)

# Title & Subtitle in single textbox
title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(4.0))
tf = title_box.text_frame
tf.word_wrap = True

# Main Title
p = tf.paragraphs[0]
p.text = "HỆ THỐNG TƯ VẤN CHỌN BẤT ĐỘNG SẢN THÔNG MINH\nTẠI THÀNH PHỐ HỒ CHÍ MINH"
p.font.name = 'Arial'
p.font.size = Pt(38)
p.font.bold = True
p.font.color.rgb = TEXT_MAIN
p.space_after = Pt(16)

# Subtitle
p2 = tf.add_paragraph()
p2.text = "Báo cáo tiến độ sơ khởi - Midterm Presentation"
p2.font.name = 'Arial'
p2.font.size = Pt(20)
p2.font.color.rgb = ACCENT_BLUE
p2.space_after = Pt(36)

# Authors/Group Info
p3 = tf.add_paragraph()
p3.text = "Thực hiện bởi: Nhóm 8 (IT2041_G8)\nĐề tài: Tư vấn Bất động sản thông minh dựa trên DSS & LLM"
p3.font.name = 'Arial'
p3.font.size = Pt(14)
p3.font.color.rgb = TEXT_MUTED

# ==============================================================================
# SLIDE 2: Input / Output Bài toán
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_title(slide, "1. Input / Output Bài toán", "Đặt tả bài toán ra quyết định lựa chọn BĐS tối ưu")

# Left Card: Input
create_card(
    slide,
    left=Inches(0.8),
    top=Inches(2.0),
    width=Inches(5.6),
    height=Inches(4.5),
    title="Đầu vào (Input)",
    items=[
        "1. Hồ sơ mong muốn (User Profile / Preferences):",
        "   - Ngân sách tối đa (Ràng buộc cứng)",
        "   - Số phòng ngủ tối thiểu (Ràng buộc cứng)",
        "   - Độ quan trọng của các tiện ích (Trọng số mềm [0, 1]): trường học, công viên, siêu thị, bệnh viện, trục đường lớn.",
        "2. Câu lệnh tự do của người dùng (User Query):",
        "   - Ví dụ: 'Tôi muốn tìm nhà yên tĩnh cho gia đình có con nhỏ, tiện đi chợ và có chỗ chơi cho bé'."
    ],
    title_color=ACCENT_BLUE
)

# Right Card: Output & Example
create_card(
    slide,
    left=Inches(6.8),
    top=Inches(2.0),
    width=Inches(5.7),
    height=Inches(4.5),
    title="Đầu ra (Output) & Ví dụ",
    items=[
        "Đầu ra (Output):",
        "  - Danh sách Top 5 Bất động sản phù hợp nhất kèm theo điểm số tổng hợp chi tiết.",
        "  - Báo cáo phân tích (Breakdown) lý do đề xuất.",
        "Ví dụ thực tế (Scenario Gia đình):",
        "  - Input: Ràng buộc: Giá <= 8 tỷ, PN >= 3. Trọng số: Trường (0.25), Công viên (0.2), Siêu thị (0.15).",
        "  - Output Đề xuất số #1: Mã căn GV_008, Điểm: 0.6558, Giá: 4.55 tỷ, 3 PN, cách trường học 614m, công viên 588m."
    ],
    title_color=ACCENT_GOLD
)

# ==============================================================================
# SLIDE 3: Pipeline cơ bản
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_title(slide, "2. Pipeline Xử lý Dữ liệu & Đề xuất", "Sơ đồ từ dữ liệu thô và yêu cầu người dùng ra kết quả Top 5")

# Draw flowchart items using card logic for clean vertical step presentation
create_card(
    slide,
    left=Inches(0.8),
    top=Inches(2.0),
    width=Inches(11.7),
    height=Inches(4.5),
    title="Quy trình xử lý tuần tự (Solution 1 Pipeline)",
    items=[
        "Bước 1: Nhận yêu cầu người dùng (Ràng buộc cứng + Trọng số ưu tiên tiện ích)",
        "Bước 2: Lọc sơ bộ (Rule-based Filtering) loại bỏ các căn vượt ngân sách hoặc thiếu số phòng ngủ mong muốn",
        "Bước 3: Làm giàu dữ liệu địa lý (POI Enrichment) tính toán khoảng cách thực tế (Haversine) từ tọa độ BĐS đến các POI trong khu vực",
        "Bước 4: Chuẩn hóa dữ liệu (Normalization) đưa các thuộc tính (khoảng cách, giá, diện tích) về thang đo [0, 1] theo xu hướng tốt hơn",
        "Bước 5: Tính điểm tổng hợp (Rule-based Scoring) nhân giá trị chuẩn hóa với trọng số của kịch bản",
        "Bước 6: Xếp hạng và xuất danh sách Top 5 đề xuất chi tiết giải thích cho người dùng"
    ],
    title_color=ACCENT_BLUE
)

# ==============================================================================
# SLIDE 4: Giải pháp đề xuất Solution 2 (Tích hợp LLM & Real-time API)
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_title(slide, "3. Giải pháp đề xuất Solution 2 (LLM & Real-time API)", "Kiến trúc hệ thống tư vấn thông minh tích hợp AI")

create_card(
    slide,
    left=Inches(0.8),
    top=Inches(2.0),
    width=Inches(5.6),
    height=Inches(4.5),
    title="Luồng xử lý thông minh",
    items=[
        "1. Semantic Query Parsing: LLM đọc hiểu mong muốn tự do bằng tiếng Việt, phân tích ra các tiêu chí ẩn (ví dụ: 'yên tĩnh' -> giảm điểm gần đường lớn).",
        "2. Real-time POI Enrichment: Sử dụng Mapbox hoặc Google Places API để tự động truy xuất các tiện ích xung quanh tọa độ của từng BĐS thực tế.",
        "3. Hybrid Scoring & Re-ranking: Kết hợp điểm số khoảng cách POI và chấm điểm nội dung bằng LLM dựa trên thông tin mô tả chi tiết của BĐS (Description).",
        "4. Tự động sinh báo cáo giải thích (LLM Explanation): Giải thích trực quan, tự nhiên và thuyết phục tại sao BĐS được chọn."
    ],
    title_color=ACCENT_BLUE
)

create_card(
    slide,
    left=Inches(6.8),
    top=Inches(2.0),
    width=Inches(5.7),
    height=Inches(4.5),
    title="Ưu điểm vượt trội so với Solution 1",
    items=[
        "🔹 Trải nghiệm người dùng cao: Không cần nhập các form khảo sát phức tạp, chỉ cần trò chuyện tự do.",
        "🔹 Độ chính xác tiện ích tuyệt đối: Lấy POI theo thời gian thực tại bất kỳ vị trí nào, không bị giới hạn bởi database tĩnh.",
        "🔹 Đánh giá BĐS sâu sắc hơn: Đọc hiểu được thông tin phi cấu trúc (ví dụ: hẻm xe hơi, nội thất cao cấp, nở hậu) từ Description tin đăng.",
        "🔹 Cá nhân hóa tối đa: Lý giải lý do đề xuất riêng biệt cho từng người dùng bằng văn phong tiếng Việt tự nhiên."
    ],
    title_color=ACCENT_GOLD
)

# ==============================================================================
# SLIDE 5: Thông tin về bộ dữ liệu
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_title(slide, "4. Thông tin Bộ dữ liệu", "Chi tiết về dữ liệu nguồn và dữ liệu làm giàu tiện ích")

create_card(
    slide,
    left=Inches(0.8),
    top=Inches(2.0),
    width=Inches(5.6),
    height=Inches(4.5),
    title="Dữ liệu Bất động sản nguồn",
    items=[
        "Bộ dữ liệu chính (HCMC Real Estate Data 2025):",
        "  - Số mẫu ban đầu: 51,304 tin đăng tại TP.HCM.",
        "  - Làm sạch & Lọc subset (Gò Vấp): Lấy 37 mẫu Nhà riêng sạch.",
        "  - Tiêu chí lọc sạch: Có tọa độ lat/lon, đầy đủ số phòng ngủ, phòng vệ sinh, giá trị hợp lý (500tr - 30 tỷ), diện tích 20 - 500m².",
        "  - Phân bố giá trị trong Gò Vấp subset: 2.85 tỷ - 27.0 tỷ.",
        "  - Diện tích subset: 26m² - 258m²."
    ],
    title_color=ACCENT_BLUE
)

create_card(
    slide,
    left=Inches(6.8),
    top=Inches(2.0),
    width=Inches(5.7),
    height=Inches(4.5),
    title="Dữ liệu tiện ích bổ sung (POIs)",
    items=[
        "Bổ sung tọa độ địa lý thực tế của các tiện ích tại Quận Gò Vấp:",
        "  - 8 Trường học nổi tiếng (Nguyễn Du, Lương Thế Vinh, Trần Hưng Đạo...)",
        "  - 4 Công viên lớn (Làng Hoa, Gia Định, phường 12...)",
        "  - 3 Bệnh viện chính (Bệnh viện Gò Vấp, Bệnh viện 175...)",
        "  - 5 Siêu thị tiện ích (Emart Gò Vấp, Coopmart Quang Trung, VinMart...)",
        "  - 4 Trục đường lớn di chuyển (Quang Trung, Nguyễn Oanh, Phạm Văn Đồng...)"
    ],
    title_color=ACCENT_GOLD
)

# ==============================================================================
# SLIDE 6: Cách đánh giá (Evaluation)
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_title(slide, "5. Cách đánh giá Hệ thống tư vấn (DSS Evaluation)", "Các tiêu chí đảm bảo chất lượng thuật toán ra quyết định")

create_card(
    slide,
    left=Inches(0.8),
    top=Inches(2.0),
    width=Inches(11.7),
    height=Inches(4.5),
    title="Tiêu chí đánh giá hệ thống tư vấn",
    items=[
        "1. Tỷ lệ thỏa mãn ràng buộc cứng (Constraint Satisfaction Rate): Cam kết 100% kết quả xuất ra không vi phạm ngân sách tối đa hoặc số phòng ngủ tối thiểu.",
        "2. Tính minh bạch điểm số (Score Transparency): Khách hàng có thể nhìn thấy chi tiết điểm thành phần (giá chuẩn hóa, khoảng cách chuẩn hóa và mức đóng góp thực tế).",
        "3. Sự tương thích theo Kịch bản (Persona Relevance): Danh sách Top 5 đề xuất phải thay đổi rõ rệt và phản ánh đúng trọng số ưu tiên của từng Persona.",
        "4. Tính ổn định xếp hạng (Ranking Stability): Thuật toán chạy deterministic, phản hồi kết quả nhất quán với cùng một cấu hình đầu vào.",
        "5. Trải nghiệm người dùng (UX Explanation): Đánh giá chất lượng của mô tả lý giải giải thích đề xuất Top 1."
    ],
    title_color=ACCENT_BLUE
)

# ==============================================================================
# SLIDE 7: Tập kiểm định & Đánh giá định lượng (Validation Suite)
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_title(slide, "6. Kiểm định Định lượng (Validation Suite)", "Thiết lập tập kiểm thử và đo lường các chỉ số ra quyết định")

create_card(
    slide,
    left=Inches(0.8),
    top=Inches(2.0),
    width=Inches(5.6),
    height=Inches(4.5),
    title="Bộ dữ liệu kiểm định (Validation Dataset)",
    items=[
        "🔹 Tập kịch bản kiểm thử (validation_scenarios.json):",
        "  - Thiết kế 5 scenarios có thuộc tính ràng buộc/mong muốn khác biệt rõ rệt.",
        "  - Định nghĩa sẵn tập nhãn BĐS tối ưu (Ground-truth Top 5) cho từng scenario làm mốc so sánh.",
        "🔹 Các chỉ số đo lường chất lượng:",
        "  - CSR (Constraint Satisfaction Rate): Tỷ lệ BĐS đạt chuẩn lọc cứng.",
        "  - Precision@5: Tỷ lệ BĐS đề xuất trùng khớp với Ground-truth.",
        "  - NDCG@5 (Ranking Quality): Chất lượng sắp xếp thứ tự ưu tiên."
    ],
    title_color=ACCENT_BLUE
)

create_card(
    slide,
    left=Inches(6.8),
    top=Inches(2.0),
    width=Inches(5.7),
    height=Inches(4.5),
    title="Kết quả kiểm tra định lượng",
    items=[
        "📊 Chỉ số đạt được trên toàn tập kiểm thử:",
        "  - CSR trung bình: 100% (Không xảy ra lỗi đề xuất vi phạm điều kiện lọc cứng).",
        "  - Precision@5 trung bình: 100% (Đề xuất chính xác toàn bộ 5 BĐS thuộc tập tối ưu).",
        "  - NDCG@5 trung bình: 1.0000 (Xếp hạng chính xác tuyệt đối theo thứ tự mong muốn).",
        "🔹 Báo cáo chi tiết định lượng đã được lưu tự động tại:",
        "  - outputs/validation_report.md"
    ],
    title_color=ACCENT_GOLD
)

# ==============================================================================
# SLIDE 8: Kết quả sơ khởi - Persona Gia đình & Người trẻ
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_title(slide, "7. Kết quả Sơ khởi - Persona 1 & 2", "Kết quả chạy thực tế của thuật toán trên subset Gò Vấp")

create_card(
    slide,
    left=Inches(0.8),
    top=Inches(2.0),
    width=Inches(5.6),
    height=Inches(4.5),
    title="Persona 1: Gia đình có con nhỏ",
    items=[
        "Ràng buộc: Giá <= 8 tỷ, PN >= 3. Ưu tiên: Trường, CV.",
        "Số căn thỏa mãn lọc: 13 / 37.",
        "Đề xuất tối ưu nhất: GV_008 (Điểm: 0.6558)",
        "  - Giá: 4.55 tỷ | Diện tích: 60m² | 3 PN - 3 WC.",
        "  - Khoảng cách tiện ích: Trường (614m), Công viên (588m), Siêu thị (186m).",
        "  - Đánh giá: Rất thích hợp cho sự phát triển của trẻ nhỏ."
    ],
    title_color=ACCENT_BLUE
)

create_card(
    slide,
    left=Inches(6.8),
    top=Inches(2.0),
    width=Inches(5.7),
    height=Inches(4.5),
    title="Persona 2: Người trẻ độc thân",
    items=[
        "Ràng buộc: Giá <= 5.5 tỷ, PN >= 2. Ưu tiên: Giá rẻ, Siêu thị.",
        "Số căn thỏa mãn lọc: 14 / 37.",
        "Đề xuất tối ưu nhất: GV_037 (Điểm: 0.6380)",
        "  - Giá: 2.85 tỷ (Rẻ nhất nhóm lọc) | Diện tích: 26m² | 2 PN - 1 WC.",
        "  - Khoảng cách tiện ích: Siêu thị (269m), Trục chính Quang Trung (472m).",
        "  - Đánh giá: Tiết kiệm ngân sách tối đa, dễ dàng di chuyển."
    ],
    title_color=ACCENT_GOLD
)

# ==============================================================================
# SLIDE 9: Kết quả sơ khởi - Persona Nhà đầu tư & Kế hoạch tiếp theo
# ==============================================================================
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_title(slide, "8. Kết quả Sơ khởi - Persona 3 & Kế hoạch", "Nhà đầu tư BĐS và các bước hoàn thiện dự án")

create_card(
    slide,
    left=Inches(0.8),
    top=Inches(2.0),
    width=Inches(5.6),
    height=Inches(4.5),
    title="Persona 3: Nhà đầu tư",
    items=[
        "Ràng buộc: Giá <= 15 tỷ, PN >= 1. Ưu tiên: Đơn giá/m² rẻ.",
        "Số căn thỏa mãn lọc: 31 / 37.",
        "Đề xuất tối ưu nhất: GV_008 (Điểm: 0.5875)",
        "  - Giá: 4.55 tỷ | Đơn giá: 75.8 triệu/m² (Rẻ nhất nhóm lọc).",
        "  - Tiện ích: Cách trục chính 727m, siêu thị 186m.",
        "  - Đánh giá: Khả năng sinh lời cao nhờ đơn giá mua vào cực tốt."
    ],
    title_color=ACCENT_BLUE
)

create_card(
    slide,
    left=Inches(6.8),
    top=Inches(2.0),
    width=Inches(5.7),
    height=Inches(4.5),
    title="Kế hoạch phát triển tiếp theo (Solution 2 Pipeline)",
    items=[
        "1. Tự động hóa lấy POI: Sử dụng Google Places API / OpenStreetMap thay vì hardcode tọa độ.",
        "2. Mở rộng bộ dữ liệu: Nâng số lượng căn lên 200 - 500 mẫu tại nhiều quận khác nhau.",
        "3. Tích hợp Mô hình ngôn ngữ lớn (LLM):",
        "   - Nhận diện yêu cầu tự do của người dùng thành tham số.",
        "   - Re-ranking dựa trên mô tả chi tiết của BĐS.",
        "   - Tự động sinh báo cáo giải thích gợi ý tự nhiên."
    ],
    title_color=ACCENT_GOLD
)

# Save the presentation
output_dir = "outputs"
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "midterm_presentation.pptx")
prs.save(output_path)
print(f"Presentation saved successfully to: {output_path}")
